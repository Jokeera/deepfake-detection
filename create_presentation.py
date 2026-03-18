#!/usr/bin/env python3
"""
Генерация презентации для защиты ВКР.
Тема: Детекция дипфейков на основе пространственно-временного анализа.

Цветовая схема: Midnight Executive
  Primary: #1E2761 (navy)
  Secondary: #CADCFC (ice blue)
  Accent: #FFFFFF (white)
  Dark accent: #408EC6 (blue)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Colors ───
NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)
BLUE = RGBColor(0x40, 0x8E, 0xC6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xFB)
DARK_TEXT = RGBColor(0x1E, 0x27, 0x61)
GRAY = RGBColor(0x6B, 0x70, 0x80)
GREEN = RGBColor(0x10, 0xB9, 0x81)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buNone = pPr.find(qn('a:buNone'))
        if buNone is not None:
            pPr.remove(buNone)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '▸'})
        pPr.append(buChar)
    return txBox


prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# =====================================================================
# SLIDE 1: Title
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, NAVY)

# Decorative bar
add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, BLUE)

add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10), Inches(1.5),
             "ДЕТЕКЦИЯ ДИПФЕЙКОВ",
             font_size=44, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.2), Inches(2.8), Inches(10), Inches(1.2),
             "на основе пространственно-временного\nанализа видеопоследовательностей",
             font_size=28, color=ICE_BLUE, bold=False, align=PP_ALIGN.LEFT)

# Divider line
add_rect(slide, Inches(1.2), Inches(4.3), Inches(3), Inches(0.04), BLUE)

add_text_box(slide, Inches(1.2), Inches(4.7), Inches(8), Inches(0.5),
             "Магистерская диссертация (ВКР)",
             font_size=18, color=ICE_BLUE, align=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.2), Inches(5.4), Inches(8), Inches(0.5),
             "НИЯУ МИФИ  ·  2026",
             font_size=16, color=GRAY, align=PP_ALIGN.LEFT)

# =====================================================================
# SLIDE 2: Актуальность
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "АКТУАЛЬНОСТЬ ИССЛЕДОВАНИЯ", font_size=28, color=WHITE, bold=True)

# Stats cards
stats = [
    ("500%", "рост числа дипфейков\nза 2019-2024 гг."),
    ("96%", "дипфейков —\nпорнографический контент"),
    ("$25B", "ущерб от fraud\nс использованием deepfake"),
]

for i, (num, desc) in enumerate(stats):
    x = Inches(0.8 + i * 4.0)
    card = add_rect(slide, x, Inches(1.5), Inches(3.5), Inches(2.0), WHITE)
    card.shadow.inherit = False
    add_rect(slide, x, Inches(1.5), Inches(3.5), Inches(0.06), BLUE)
    add_text_box(slide, x + Inches(0.3), Inches(1.7), Inches(3), Inches(0.8),
                 num, font_size=40, color=BLUE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3), Inches(0.8),
                 desc, font_size=14, color=GRAY)

add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(11), Inches(3.0), [
    "Генеративные модели (GAN, diffusion) делают фальсификации неотличимыми для человека",
    "Угрозы: политические манипуляции, финансовое мошенничество, нарушение приватности",
    "Необходимы автоматические методы детекции, учитывающие как пространственные, так и временные артефакты",
], font_size=16)

# =====================================================================
# SLIDE 3: Цель и задачи
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ЦЕЛЬ И ЗАДАЧИ ИССЛЕДОВАНИЯ", font_size=28, color=WHITE, bold=True)

# Goal card
goal_card = add_rect(slide, Inches(0.8), Inches(1.4), Inches(11.5), Inches(1.3), NAVY)
add_text_box(slide, Inches(1.2), Inches(1.5), Inches(10.5), Inches(0.3),
             "ЦЕЛЬ", font_size=14, color=BLUE, bold=True)
add_text_box(slide, Inches(1.2), Inches(1.85), Inches(10.5), Inches(0.7),
             "Разработка и исследование метода детекции дипфейков на основе\n"
             "параллельного пространственно-временного анализа видеопоследовательностей",
             font_size=17, color=WHITE)

# Tasks
tasks = [
    "Исследовать существующие подходы к детекции дипфейков",
    "Спроектировать двухпоточную (dual-path) архитектуру с adaptive fusion",
    "Реализовать полный pipeline: предобработка → обучение → оценка → инференс",
    "Провести серию экспериментов (ablation study) для валидации архитектурных решений",
    "Разработать MVP-систему с веб-интерфейсом для демонстрации",
]

for i, task in enumerate(tasks):
    y = Inches(3.1 + i * 0.7)
    # Number circle
    add_rect(slide, Inches(0.8), y, Inches(0.45), Inches(0.45), BLUE)
    add_text_box(slide, Inches(0.8), y, Inches(0.45), Inches(0.45),
                 str(i + 1), font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.5), y, Inches(10), Inches(0.45),
                 task, font_size=16, color=DARK_TEXT)

# =====================================================================
# SLIDE 4: Обзор методов
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ОБЗОР СУЩЕСТВУЮЩИХ МЕТОДОВ", font_size=28, color=WHITE, bold=True)

methods = [
    ("Пространственные", "NAVY",
     ["XceptionNet, EfficientNet",
      "Анализ артефактов в отдельных кадрах",
      "Текстуры, границы, частотные аномалии",
      "❌ Не учитывают временную динамику"]),
    ("Временные", "BLUE",
     ["LSTM, 3D CNN (I3D, SlowFast)",
      "Анализ межкадровых изменений",
      "Мерцание, несогласованность движений",
      "❌ Теряют пространственные детали"]),
    ("Гибридные (наш подход)", "GREEN",
     ["Параллельные spatial + temporal ветви",
      "Совместный анализ обоих типов артефактов",
      "Adaptive fusion весов ветвей",
      "✅ Полный пространственно-временной анализ"]),
]

colors_map = {"NAVY": NAVY, "BLUE": BLUE, "GREEN": GREEN}

for i, (title, clr, items) in enumerate(methods):
    x = Inches(0.6 + i * 4.15)
    card = add_rect(slide, x, Inches(1.4), Inches(3.8), Inches(5.0), WHITE)
    add_rect(slide, x, Inches(1.4), Inches(3.8), Inches(0.06), colors_map[clr])
    add_text_box(slide, x + Inches(0.2), Inches(1.6), Inches(3.4), Inches(0.5),
                 title, font_size=18, color=colors_map[clr], bold=True)
    add_bullet_list(slide, x + Inches(0.2), Inches(2.2), Inches(3.4), Inches(3.5),
                    items, font_size=14, color=DARK_TEXT)

# =====================================================================
# SLIDE 5: Предложенный метод (Архитектура)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ПРЕДЛОЖЕННАЯ АРХИТЕКТУРА", font_size=28, color=WHITE, bold=True)

# Architecture diagram using shapes
# Spatial branch
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(2.5), WHITE)
add_rect(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.06), BLUE)
add_text_box(slide, Inches(0.7), Inches(1.65), Inches(5), Inches(0.4),
             "SPATIAL BRANCH", font_size=16, color=BLUE, bold=True)

blocks_s = ["Видеокадры\n[B,16,3,224,224]", "EfficientNet-B4\n(frozen→unfreeze)", "Mean Pool\n→ 512-dim"]
for j, txt in enumerate(blocks_s):
    bx = Inches(0.7 + j * 1.8)
    add_rect(slide, bx, Inches(2.2), Inches(1.6), Inches(1.3), ICE_BLUE)
    add_text_box(slide, bx + Inches(0.05), Inches(2.3), Inches(1.5), Inches(1.1),
                 txt, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

# Temporal branch
add_rect(slide, Inches(0.5), Inches(4.3), Inches(5.5), Inches(2.5), WHITE)
add_rect(slide, Inches(0.5), Inches(4.3), Inches(5.5), Inches(0.06), NAVY)
add_text_box(slide, Inches(0.7), Inches(4.45), Inches(5), Inches(0.4),
             "TEMPORAL BRANCH", font_size=16, color=NAVY, bold=True)

blocks_t = ["Frame Diffs\n[B,15,3,128,128]", "EfficientNet-B0\n+ Transformer", "CLS Token\n→ 512-dim"]
for j, txt in enumerate(blocks_t):
    bx = Inches(0.7 + j * 1.8)
    add_rect(slide, bx, Inches(5.0), Inches(1.6), Inches(1.3), ICE_BLUE)
    add_text_box(slide, bx + Inches(0.05), Inches(5.1), Inches(1.5), Inches(1.1),
                 txt, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

# Fusion + Head
add_rect(slide, Inches(6.8), Inches(2.5), Inches(2.5), Inches(2.0), BLUE)
add_text_box(slide, Inches(6.9), Inches(2.7), Inches(2.3), Inches(1.6),
             "Adaptive\nWeighted\nFusion\n\nα₁·h_s + α₂·h_t",
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(10.0), Inches(2.5), Inches(2.5), Inches(2.0), NAVY)
add_text_box(slide, Inches(10.1), Inches(2.7), Inches(2.3), Inches(1.6),
             "Classification\nHead\n\n512→256→128→1\nBatchNorm + Dropout",
             font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Arrows (simple rectangles as connectors)
add_rect(slide, Inches(6.1), Inches(2.6), Inches(0.6), Inches(0.06), BLUE)
add_rect(slide, Inches(6.1), Inches(5.4), Inches(0.6), Inches(0.06), NAVY)
add_rect(slide, Inches(9.4), Inches(3.4), Inches(0.5), Inches(0.06), WHITE)

# =====================================================================
# SLIDE 6: Датасет и предобработка
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ДАТАСЕТ И ПРЕДОБРАБОТКА", font_size=28, color=WHITE, bold=True)

# Dataset info cards
ds_items = [
    ("DFDC", "DeepFake Detection Challenge\n~3,300 видео"),
    ("16 кадров", "Униформная выборка\nиз каждого видео"),
    ("70/15/15", "Train / Val / Test\nVideo-level split"),
]

for i, (title, desc) in enumerate(ds_items):
    x = Inches(0.8 + i * 4.0)
    add_rect(slide, x, Inches(1.4), Inches(3.5), Inches(1.5), WHITE)
    add_rect(slide, x, Inches(1.4), Inches(0.06), Inches(1.5), BLUE)
    add_text_box(slide, x + Inches(0.3), Inches(1.5), Inches(3), Inches(0.5),
                 title, font_size=22, color=BLUE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.0), Inches(3), Inches(0.7),
                 desc, font_size=14, color=GRAY)

# Pipeline
add_text_box(slide, Inches(0.8), Inches(3.3), Inches(10), Inches(0.5),
             "Pipeline предобработки", font_size=20, color=NAVY, bold=True)

pipeline_steps = [
    "MTCNN\nFace Detection",
    "Crop + Margin\n(20%)",
    "Resize\n112×112",
    "16 Frames\nUniform Sample",
    "Clip-Consistent\nAugmentation",
]

for i, step in enumerate(pipeline_steps):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(4.0), Inches(2.2), Inches(1.2), WHITE)
    add_rect(slide, x, Inches(4.0), Inches(2.2), Inches(0.06), BLUE)
    add_text_box(slide, x + Inches(0.1), Inches(4.15), Inches(2.0), Inches(1.0),
                 step, font_size=13, color=NAVY, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5),
             "Аугментация: HorizontalFlip, Brightness/Contrast ±0.1, JPEG compression (p=0.3)\n"
             "Нормализация: ImageNet stats (spatial), per-clip z-score (temporal)",
             font_size=14, color=GRAY)

# =====================================================================
# SLIDE 7: Эксперименты (Ablation Study)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ЭКСПЕРИМЕНТЫ: ABLATION STUDY", font_size=28, color=WHITE, bold=True)

# Table header
table_data = [
    ["Эксперимент", "Модель", "AUC ↑", "Accuracy", "F1", "EER ↓"],
    ["A1 (proposed)", "Dual-Path + Adaptive Fusion", "—", "—", "—", "—"],
    ["A2", "Spatial-only (EfficientNet-B4)", "—", "—", "—", "—"],
    ["A3", "Temporal-only (EfficientNet-B0 + Transformer)", "—", "—", "—", "—"],
    ["A4", "Sequential (CNN → BiLSTM)", "—", "—", "—", "—"],
]

rows = len(table_data)
cols = len(table_data[0])
table_shape = slide.shapes.add_table(rows, cols,
                                      Inches(0.5), Inches(1.4),
                                      Inches(12.3), Inches(3.0))
table = table_shape.table

col_widths = [Inches(2.0), Inches(4.5), Inches(1.3), Inches(1.5), Inches(1.3), Inches(1.3)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = table.cell(r, c)
        cell.text = table_data[r][c]
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER if c >= 2 else PP_ALIGN.LEFT
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            else:
                paragraph.font.color.rgb = DARK_TEXT
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif r == 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ICE_BLUE
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(2.0),
             "⚠ Результаты будут заполнены после завершения 30-эпохального обучения.\n"
             "Метрики вычисляются на held-out test set (15% данных, video-level split).\n\n"
             "Гиперпараметры: batch_size=8, lr_head=3e-4, lr_backbone=1e-4,\n"
             "warmup=5 эпох (frozen backbone), patience=5 (early stopping по AUC)",
             font_size=14, color=GRAY)

# =====================================================================
# SLIDE 8: Результаты (placeholder для графиков)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "РЕЗУЛЬТАТЫ ЭКСПЕРИМЕНТОВ", font_size=28, color=WHITE, bold=True)

# 4 placeholder boxes for charts
chart_titles = ["Loss Curves", "Validation AUC", "ROC Curve", "Confusion Matrix"]
for i, title in enumerate(chart_titles):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 3.0)
    card = add_rect(slide, x, y, Inches(5.8), Inches(2.6), WHITE)
    add_rect(slide, x, y, Inches(5.8), Inches(0.05), BLUE)
    add_text_box(slide, x, y + Inches(0.8), Inches(5.8), Inches(1.0),
                 f"[ {title} ]\nВставить график после экспериментов",
                 font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 9: MVP система
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "MVP: ВЕБ-ПРИЛОЖЕНИЕ", font_size=28, color=WHITE, bold=True)

# Left: features
features = [
    "Flask веб-приложение (порт 7860)",
    "Загрузка видео через браузер (до 300 MB)",
    "Автоматическое извлечение лиц (MTCNN)",
    "Video-level inference с визуализацией",
    "JSON-результат с probability и fusion weights",
    "Мультиязычный интерфейс (RU/EN)",
    "Docker-контейнеризация для деплоя",
]

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.5),
             "Возможности системы", font_size=20, color=NAVY, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5),
                features, font_size=15, color=DARK_TEXT)

# Right: tech stack card
add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5), WHITE)
add_rect(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.06), BLUE)
add_text_box(slide, Inches(7.3), Inches(1.5), Inches(5), Inches(0.5),
             "Стек технологий", font_size=18, color=BLUE, bold=True)

stack = [
    "Backend:  Python 3.12, Flask, PyTorch",
    "Models:   EfficientNet (timm), Transformer",
    "Face Det: MTCNN (facenet-pytorch)",
    "Frontend: HTML/CSS/JS (inline)",
    "Deploy:   Docker, python:3.12-slim",
    "Device:   CPU / CUDA / MPS (Apple Silicon)",
]
add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5), Inches(4.0),
                stack, font_size=14, color=DARK_TEXT)

add_text_box(slide, Inches(7.3), Inches(5.5), Inches(5), Inches(1.0),
             "[ Скриншот MVP ]\nВставить скриншот веб-интерфейса",
             font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 10: Заключение
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, LIGHT_BG)

add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ЗАКЛЮЧЕНИЕ", font_size=28, color=WHITE, bold=True)

# Two columns
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.5),
             "Основные результаты", font_size=20, color=NAVY, bold=True)

results_items = [
    "Разработана dual-path архитектура для пространственно-временного анализа",
    "Реализован полный pipeline от предобработки до инференса",
    "Проведён ablation study (4 эксперимента) для валидации решений",
    "Создано MVP веб-приложение с Docker-контейнеризацией",
]
add_bullet_list(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(3.5),
                results_items, font_size=15, color=DARK_TEXT)

add_text_box(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(0.5),
             "Направления развития", font_size=20, color=NAVY, bold=True)

future_items = [
    "Расширение на мультимодальный анализ (аудио + видео)",
    "Attention-based pooling вместо mean pooling",
    "Cross-dataset evaluation (CelebDF, FaceForensics++)",
    "Real-time processing для видеопотоков",
]
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.5), Inches(3.5),
                future_items, font_size=15, color=DARK_TEXT)

# =====================================================================
# SLIDE 11: Спасибо за внимание
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, NAVY)

add_rect(slide, Inches(0), Inches(0), Inches(0.15), SLIDE_H, BLUE)

add_text_box(slide, Inches(1.2), Inches(2.0), Inches(10), Inches(1.5),
             "СПАСИБО ЗА ВНИМАНИЕ!",
             font_size=44, color=WHITE, bold=True, align=PP_ALIGN.LEFT)

add_rect(slide, Inches(1.2), Inches(3.8), Inches(3), Inches(0.04), BLUE)

add_text_box(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(1.0),
             "Готов ответить на вопросы",
             font_size=24, color=ICE_BLUE, align=PP_ALIGN.LEFT)

add_text_box(slide, Inches(1.2), Inches(5.5), Inches(8), Inches(0.5),
             "НИЯУ МИФИ  ·  2026",
             font_size=16, color=GRAY, align=PP_ALIGN.LEFT)

# ─── Save ───
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "presentation_defense.pptx"
)
prs.save(output_path)
print(f"Презентация сохранена: {output_path}")
